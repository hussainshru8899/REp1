import os
import pandas as pd
from pathlib import Path
from utils.ocr_utils import ocr_pdf_to_text, ocr_image_to_text
from pdf2image import convert_from_path


def convert_file_dispatch(in_path, tool, tmpdir):
    ext = Path(in_path).suffix.lower()

    if tool == 'pdf_to_excel' and ext == '.pdf':
        return pdf_to_excel(in_path, tmpdir)
    if tool == 'excel_to_pdf' and ext in ['.xls', '.xlsx']:
        return excel_to_pdf(in_path, tmpdir)
    if tool == 'jpg_to_excel' and ext in ['.jpg', '.jpeg', '.png']:
        return image_to_excel(in_path, tmpdir)
    if tool == 'img_convert' and ext in ['.jpg', '.jpeg', '.png']:
        return image_format_convert(in_path, tmpdir)
    if tool == 'pptx_to_xlsx' and ext in ['.pptx']:
        return pptx_to_xlsx(in_path, tmpdir)
    if tool == 'xlsx_to_pptx' and ext in ['.xlsx']:
        return xlsx_to_pptx(in_path, tmpdir)

    return None


# --- Implementations (basic/easy-to-extend) ---

def pdf_to_excel(pdf_path, out_dir):
    # Try Camelot
    try:
        import camelot
        tables = camelot.read_pdf(pdf_path, pages='all')
        if tables and len(tables) > 0:
            out_path = os.path.join(out_dir, Path(pdf_path).stem + '.xlsx')
            with pd.ExcelWriter(out_path, engine='openpyxl') as writer:
                for i, t in enumerate(tables, start=1):
                    df = t.df
                    df.to_excel(writer, sheet_name=f'sheet{i}', index=False)
            return out_path
    except Exception as e:
        print('camelot error', e)

    # Fallback: OCR
    text_pages = ocr_pdf_to_text(pdf_path)
    out_path = os.path.join(out_dir, Path(pdf_path).stem + '_ocr.xlsx')
    df = pd.DataFrame({'page_text': text_pages})
    df.to_excel(out_path, index=False)
    return out_path


def excel_to_pdf(xlsx_path, out_dir):
    # Simple rendering: convert each sheet to an HTML table and to PDF using wkhtmltopdf/poppler (poppler-utils needed)
    try:
        import pdfkit
    except Exception:
        pdfkit = None

    x = pd.ExcelFile(xlsx_path)
    html = ''
    for sheet in x.sheet_names:
        df = x.parse(sheet)
        html += f'<h2>{sheet}</h2>' + df.to_html(index=False)

    out_pdf = os.path.join(out_dir, Path(xlsx_path).stem + '.pdf')
    if pdfkit:
        pdfkit.from_string(html, out_pdf)
        return out_pdf
    else:
        # fallback to simple text PDF via pillow
        from PIL import Image, ImageDraw, ImageFont
        img = Image.new('RGB', (2480, 3508), color='white')
        d = ImageDraw.Draw(img)
        d.text((100,100), html[:8000], fill=(0,0,0))
        img.save(out_pdf, 'PDF', resolution=100.0)
        return out_pdf


def image_to_excel(img_path, out_dir):
    # OCR image and save text to Excel
    text = ocr_image_to_text(img_path)
    import pandas as pd
    out_path = os.path.join(out_dir, Path(img_path).stem + '.xlsx')
    df = pd.DataFrame({'text': [text]})
    df.to_excel(out_path, index=False)
    return out_path


def image_format_convert(img_path, out_dir):
    from PIL import Image
    im = Image.open(img_path)
    base = Path(img_path).stem
    # convert png->jpg or jpg->png
    if img_path.lower().endswith('.png'):
        out_path = os.path.join(out_dir, base + '.jpg')
        rgb = im.convert('RGB')
        rgb.save(out_path, 'JPEG')
    else:
        out_path = os.path.join(out_dir, base + '.png')
        im.save(out_path, 'PNG')
    return out_path


def pptx_to_xlsx(pptx_path, out_dir):
    # Basic: extract text from slides and write to excel
    from pptx import Presentation
    prs = Presentation(pptx_path)
    rows = []
    for i, slide in enumerate(prs.slides, start=1):
        text = '\n'.join([shape.text for shape in slide.shapes if hasattr(shape, 'text')])
        rows.append({'slide': i, 'text': text})
    import pandas as pd
    out_path = os.path.join(out_dir, Path(pptx_path).stem + '.xlsx')
    pd.DataFrame(rows).to_excel(out_path, index=False)
    return out_path


def xlsx_to_pptx(xlsx_path, out_dir):
    # Basic: render each sheet as table text on a slide
    from pptx import Presentation
    prs = Presentation()
    x = pd.ExcelFile(xlsx_path)
    for sheet in x.sheet_names:
        df = x.parse(sheet)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        text = df.to_string(index=False)
        from pptx.util import Inches
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = text[:10000]
    out_path = os.path.join(out_dir, Path(xlsx_path).stem + '.pptx')
    prs.save(out_path)
    return out_path
